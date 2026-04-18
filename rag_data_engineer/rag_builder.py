"""Document parsing, chunking, embedding and Chroma persistence."""
from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from typing import Iterable

import chromadb
import tiktoken
import voyageai
from chromadb.config import Settings as ChromaSettings
from pypdf import PdfReader

from .drive_loader import DriveFile
from .obsidian import preprocess_obsidian

logger = logging.getLogger(__name__)


def make_chroma_client(
    *,
    persist_dir: str | None = None,
    http_host: str | None = None,
    http_port: int = 8000,
    http_ssl: bool = False,
    auth_token: str | None = None,
):
    """Return a Chroma client. HTTP mode wins when http_host is set."""
    if http_host:
        headers = {}
        if auth_token:
            headers["Authorization"] = f"Bearer {auth_token}"
        return chromadb.HttpClient(
            host=http_host,
            port=http_port,
            ssl=http_ssl,
            headers=headers or None,
            settings=ChromaSettings(anonymized_telemetry=False),
        )
    if not persist_dir:
        raise ValueError("Either http_host or persist_dir must be provided")
    return chromadb.PersistentClient(
        path=persist_dir,
        settings=ChromaSettings(anonymized_telemetry=False),
    )


@dataclass
class Chunk:
    text: str
    metadata: dict
    chunk_id: str


def extract_text(file: DriveFile) -> tuple[str, dict]:
    """Extract text from a DriveFile. Returns (text, extra_metadata).

    For Obsidian markdown files, frontmatter tags/aliases are extracted
    and wikilink syntax is cleaned up.
    """
    extra_meta: dict = {}
    mime = file.effective_mime
    if mime == "application/pdf":
        if not file.content:
            return "", extra_meta
        reader = PdfReader(io.BytesIO(file.content))
        return "\n".join((page.extract_text() or "") for page in reader.pages), extra_meta
    if mime in {"text/plain", "text/markdown", "text/csv"}:
        raw = file.content.decode("utf-8", errors="replace")
        if file.name.endswith(".md"):
            raw, extra_meta = preprocess_obsidian(raw)
        return raw, extra_meta
    if (
        mime
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        from docx import Document  # lazy import

        doc = Document(io.BytesIO(file.content))
        return "\n".join(p.text for p in doc.paragraphs), extra_meta
    if (
        mime
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        from pptx import Presentation  # lazy import

        prs = Presentation(io.BytesIO(file.content))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts), extra_meta
    logger.warning("No extractor for mime %s (file %s)", mime, file.name)
    return "", extra_meta


def chunk_text(text: str, chunk_size: int, overlap: int, encoder) -> list[str]:
    """Token-based chunking with overlap. Stable and embedding-friendly."""
    if not text.strip():
        return []
    tokens = encoder.encode(text)
    chunks: list[str] = []
    step = max(chunk_size - overlap, 1)
    for start in range(0, len(tokens), step):
        window = tokens[start : start + chunk_size]
        if not window:
            break
        chunks.append(encoder.decode(window))
        if start + chunk_size >= len(tokens):
            break
    return chunks


class RagBuilder:
    def __init__(
        self,
        voyage_api_key: str,
        embedding_model: str,
        collection_name: str,
        chunk_size: int,
        chunk_overlap: int,
        *,
        persist_dir: str | None = None,
        http_host: str | None = None,
        http_port: int = 8000,
        http_ssl: bool = False,
        auth_token: str | None = None,
    ):
        self.client = voyageai.Client(api_key=voyage_api_key)
        self.embedding_model = embedding_model
        self.chroma = make_chroma_client(
            persist_dir=persist_dir,
            http_host=http_host,
            http_port=http_port,
            http_ssl=http_ssl,
            auth_token=auth_token,
        )
        self.collection = self.chroma.get_or_create_collection(
            name=collection_name, metadata={"hnsw:space": "cosine"}
        )
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        # Voyage doesn't expose tiktoken-compatible tokenizers; cl100k_base
        # is a close-enough approximation for chunk-size budgeting.
        self.encoder = tiktoken.get_encoding("cl100k_base")

    # ---- state helpers --------------------------------------------------

    def known_files(self) -> dict[str, str]:
        """Return {file_id: modified_time} already indexed."""
        # We store one Chroma metadata 'file_id' + 'modified_time' per chunk.
        # A single get() returns all metadata.
        out: dict[str, str] = {}
        offset = 0
        page = 1000
        while True:
            res = self.collection.get(limit=page, offset=offset, include=["metadatas"])
            metas = res.get("metadatas") or []
            if not metas:
                break
            for m in metas:
                fid = m.get("file_id")
                mt = m.get("modified_time")
                if fid and mt and out.get(fid, "") < mt:
                    out[fid] = mt
            if len(metas) < page:
                break
            offset += page
        return out

    def delete_file(self, file_id: str) -> None:
        self.collection.delete(where={"file_id": file_id})

    # ---- ingestion ------------------------------------------------------

    def _embed(self, texts: list[str]) -> list[list[float]]:
        # Voyage allows up to 128 inputs per request for voyage-3 family.
        out: list[list[float]] = []
        for i in range(0, len(texts), 128):
            batch = texts[i : i + 128]
            resp = self.client.embed(
                texts=batch,
                model=self.embedding_model,
                input_type="document",
            )
            out.extend(resp.embeddings)
        return out

    def upsert_file(self, file: DriveFile) -> int:
        text, extra_meta = extract_text(file)
        chunks = chunk_text(
            text, self.chunk_size, self.chunk_overlap, self.encoder
        )
        if not chunks:
            return 0

        # Replace any prior version of this file.
        self.delete_file(file.id)

        ids = [f"{file.id}::{i}" for i in range(len(chunks))]
        metadatas = [
            {
                "file_id": file.id,
                "file_name": file.name,
                "mime_type": file.mime_type,
                "modified_time": file.modified_time,
                "chunk_index": i,
                **extra_meta,
            }
            for i in range(len(chunks))
        ]
        embeddings = self._embed(chunks)
        self.collection.upsert(
            ids=ids, documents=chunks, metadatas=metadatas, embeddings=embeddings
        )
        return len(chunks)

    def reconcile(
        self, files: Iterable[DriveFile], known: dict[str, str]
    ) -> dict[str, int]:
        """Add new / update modified / leave untouched. Returns counters."""
        stats = {"added": 0, "updated": 0, "skipped": 0, "chunks": 0, "failed": 0}
        seen_ids: set[str] = set()
        for f in files:
            seen_ids.add(f.id)
            prev = known.get(f.id)
            if prev == f.modified_time:
                stats["skipped"] += 1
                continue
            try:
                n = self.upsert_file(f)
            except Exception as exc:  # noqa: BLE001 — never fail the whole batch
                logger.warning(
                    "Failed to ingest %s (%s): %s", f.name, f.id, exc
                )
                stats["failed"] += 1
                continue
            stats["chunks"] += n
            if prev is None:
                stats["added"] += 1
            else:
                stats["updated"] += 1
        # Delete files removed from Drive
        deleted = 0
        for fid in known.keys() - seen_ids:
            self.delete_file(fid)
            deleted += 1
        stats["deleted"] = deleted
        return stats
